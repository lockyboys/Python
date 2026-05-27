# [document_analyzer.py] - 문서 분석 모듈 (v17)
import os
import config
import utils
from pathlib import Path

# -----------------------------------------
# [개선] 라이브러리 로드 실패 시에도 시스템이 계속 작동하도록 예외 처리 강화
# [개선] 각 라이브러리 로드 시 상세 오류 로그 기록
# [개선] 문서 분석에 필요한 라이브러리 로드 시도 및 상태 플래그 설정 
# -----------------------------------------
try:
    import fitz
    config.PDF_READY = True
except Exception as e: 
    utils.log_error(f"PDF 라이브러리(fitz) 로드 실패: {e}", False)
    pass

try:
    from docx import Document
    config.DOCX_READY = True
except Exception as e: 
    utils.log_error(f"Word 라이브러리(docx) 로드 실패: {e}", False)
    pass

try:
    import openpyxl
    config.EXCEL_READY = True
except Exception as e: 
    utils.log_error(f"Excel 라이브러리(openpyxl) 로드 실패: {e}", False)
    pass

try:
    from pptx import Presentation
    config.PPT_READY = True
except Exception as e: 
    utils.log_error(f"PPT 라이브러리(pptx) 로드 실패: {e}", False)
    pass

try:
    import olefile
    config.HWP_READY = True
except Exception as e: 
    utils.log_error(f"HWP 라이브러리(olefile) 로드 실패: {e}", False)
    pass
# ----------------------------------------
# [핵심] 문서 내용 분석 및 분류 함수
# [개선] 분석 실패 시에도 기본 분류로 이동하도록 보완
# [개선] 분석 과정에서 제외 폴더 내 파일 보호 로직 추가 (해체 모드가 아닐 때)
# [개선] 분석 중 발생하는 오류는 로그에 기록하되, 시스템이 계속 작동하도록 예외 처리 강화
# [개선] 분석된 문서 수 카운트 및 반환
# [개선] 분석 시작 및 완료 로그 추가
# [개선] 분석 중인 파일명 로그 추가 (선택적)
# ----------------------------------------
def extract_text(file_path):
    ext = file_path.suffix.lower()
    text = ""
    try:
        if ext == ".txt":
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read(config.MAX_DOC_TEXT_LENGTH)
        elif ext == ".pdf" and PDF_READY:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text()
                if len(text) > config.MAX_DOC_TEXT_LENGTH: break
            doc.close()
        elif ext in [".docx", ".doc"] and DOCX_READY:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs[:50]])
        elif ext in [".xlsx", ".xls"] and EXCEL_READY:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            text = " ".join([str(c) for row in wb.active.iter_rows(max_row=50, values_only=True) for c in row if c])
        elif ext in [".pptx", ".ppt"] and PPT_READY:
            prs = Presentation(file_path)
            for slide in prs.slides[:10]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + " "
        elif ext in [".hwp", ".hwpx"] and HWP_READY:
            f = olefile.OleFileIO(file_path)
            if 'PrvText' in f.listdir():
                text = f.openstream('PrvText').read().decode('utf-16', errors='ignore')
            f.close()
    except Exception as e:
        utils.log_error(f"문서 텍스트 추출 오류 ({file_path.name}): {e}")
    return text[:config.MAX_DOC_TEXT_LENGTH]
# ----------------------------------------
# [핵심] 문서 분석 및 분류 로직
# [개선] 분석 실패 시에도 기본 분류로 이동하도록 보완
# [개선] 분석 과정에서 제외 폴더 내 파일 보호 로직 추가 (해체 모드가 아닐 때)
# [개선] 분석 중 발생하는 오류는 로그에 기록하되, 시스템이 계속 작동하도록 예외 처리 강화
# [개선] 분석된 문서 수 카운트 및 반환
# [개선] 분석 시작 및 완료 로그 추가
# [개선] 분석 중인 파일명 로그 추가 (선택적)
# ----------------------------------------

def analyze_document_content(file_path):
    try:
        # [핵심] 경로 기반 제외 체크
        if utils.is_excluded(file_path): return None

        name = file_path.name.lower()
        for folder, kws in config.KEYWORD_RULES.items():
            if any(kw.lower() in name for kw in kws): return folder
        content = extract_text(file_path).lower()
        if content:
            for folder, kws in config.KEYWORD_RULES.items():
                if any(kw.lower() in content for kw in kws): return folder
        for folder, exts in config.EXTENSION_RULES.items():
            if file_path.suffix.lower() in exts: return folder
    except Exception as e:
        utils.log_error(f"문서 내용 분석 오류 ({file_path.name}): {e}")
    return "99_미분류_기타"
#-----------------------------------------
# [핵심] 문서 분석 및 분류 실행 함수
# [개선] 예외 처리 강화 및 분석 실패 시에도 기본 분류로 이동하도록 보완
# [개선] 분석 과정에서 제외 폴더 내 파일 보호 로직 추가 (해체 모드가 아닐 때)
# [개선] 분석 중 발생하는 오류는 로그에 기록하되, 시스템이 계속 작동하도록 예외 처리 강화
# [개선] 분석된 문서 수 카운트 및 반환
# [개선] 분석 시작 및 완료 로그 추가
# [개선] 분석 중인 파일명 로그 추가 (선택적)
# [개선] 분석된 문서 수 카운트 및 반환
# [개선] 분석 시작 및 완료 로그 추가
# [개선] 분석 중인 파일명 로그 추가 (선택적)
#  ----------------------------------------
def run_document_organizing(target_path):
    target = Path(target_path)
    count = 0
    print("📄 문서 지능형 내용 분석 및 분류 중...")
    utils.log_message("📄 문서 지능형 내용 분석 및 분류 중 ...", "INFO")
    try:
        pattern = '**/*' if (config.RECURSIVE_SCAN or config.UNPACK_ALL) else '*'
        doc_files = [f for f in target.glob(pattern) if f.is_file() and f.suffix.lower() in config.DOCUMENT_EXTENSIONS]
        for item in doc_files:
            # [핵심] 경로 기반 제외 체크 (해체 모드가 아닐 때)
            if utils.is_excluded(item): continue
            if not config.UNPACK_ALL:
                if any(p.name.startswith(('0', '1', 'AI_TF', '영상_그룹')) for p in item.parents if p != target):
                    continue
            category = analyze_document_content(item)
            utils.move_file(item, target / category)
            count += 1
    except Exception as e:
        utils.log_error(f"문서 정리 프로세스 오류: {e}")
    return count
